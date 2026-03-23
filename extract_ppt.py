from pptx import Presentation

ppt_path = r'C:\Users\naray\OneDrive\Desktop\Purna\traffic-insight-dashboard\ROLL_NUMBER(Review-1).pptx'
prs = Presentation(ppt_path)

for slide_num, slide in enumerate(prs.slides, 1):
    print(f'\n========== SLIDE {slide_num} ==========')
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            print(shape.text)
